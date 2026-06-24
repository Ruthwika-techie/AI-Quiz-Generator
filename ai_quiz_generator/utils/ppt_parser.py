"""
PPT Parser Module
=================
Extracts text content from PowerPoint (.ppt/.pptx) files using python-pptx.
Handles both .pptx and legacy .ppt formats with appropriate error handling.
"""

import os
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt


class PPTParserError(Exception):
    """Custom exception for PPT parsing errors."""
    pass


def validate_ppt_file(file_path: str) -> bool:
    """
    Validate that the file is a proper PowerPoint file.
    
    Args:
        file_path: Path to the file to validate
        
    Returns:
        True if valid, raises exception otherwise
    """
    if not os.path.exists(file_path):
        raise PPTParserError(f"File not found: {file_path}")
    
    ext = Path(file_path).suffix.lower()
    if ext not in ('.ppt', '.pptx'):
        raise PPTParserError(f"Invalid file format: {ext}. Only .ppt and .pptx files are supported.")
    
    return True


def extract_ppt_content(file_path: str) -> Dict:
    """
    Extract all text content from a PowerPoint file.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        Dictionary containing:
            - slide_count: Total number of slides
            - slides: List of slide data (each with slide_number and text)
            - all_text: Combined text from all slides
            - file_name: Original file name
            - file_size: File size in bytes
    """
    validate_ppt_file(file_path)
    
    file_name = Path(file_path).name
    file_size = os.path.getsize(file_path)
    
    try:
        # Open the presentation
        prs = Presentation(file_path)
        
        slides_data = []
        all_text_parts = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
                
                # Handle tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text_parts.append(" | ".join(row_text))
                
                # Handle grouped shapes
                if shape.shape_type == 6:  # Group shape
                    try:
                        for child_shape in shape.shapes:
                            if hasattr(child_shape, "text") and child_shape.text.strip():
                                slide_text_parts.append(child_shape.text.strip())
                    except Exception:
                        pass
            
            slide_text = "\n".join(slide_text_parts)
            
            slides_data.append({
                "slide_number": idx,
                "text": slide_text,
                "word_count": len(slide_text.split()) if slide_text else 0
            })
            
            if slide_text:
                all_text_parts.append(slide_text)
        
        return {
            "slide_count": len(slides_data),
            "slides": slides_data,
            "all_text": "\n\n".join(all_text_parts),
            "file_name": file_name,
            "file_size": file_size,
            "total_word_count": sum(s["word_count"] for s in slides_data)
        }
    
    except Exception as e:
        raise PPTParserError(f"Failed to parse PowerPoint file: {str(e)}")


def get_slide_preview(slide_text: str, max_chars: int = 200) -> str:
    """
    Get a preview of slide text, truncated to max_chars.
    
    Args:
        slide_text: The full text of the slide
        max_chars: Maximum number of characters for preview
        
    Returns:
        Truncated preview string
    """
    if not slide_text:
        return "(No text content)"
    
    if len(slide_text) <= max_chars:
        return slide_text
    
    return slide_text[:max_chars].rsplit(' ', 1)[0] + "..."


def format_file_size(size_bytes: int) -> str:
    """
    Format file size in human-readable format.
    
    Args:
        size_bytes: File size in bytes
        
    Returns:
        Formatted string like "2.5 MB"
    """
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"


def extract_key_concepts(text: str) -> List[str]:
    """
    Extract key concepts from text by finding capitalized phrases and key terms.
    Used as a fallback when no AI is available.
    
    Args:
        text: The extracted text from slides
        
    Returns:
        List of key concept strings
    """
    import re
    
    concepts = set()
    
    # Find capitalized phrases (potential proper nouns, key terms)
    capitalized_phrases = re.findall(r'([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)', text)
    for phrase in capitalized_phrases:
        if 3 <= len(phrase) <= 50:
            concepts.add(phrase)
    
    # Find words in quotes
    quoted_terms = re.findall(r'"([^"]+)"', text)
    for term in quoted_terms:
        if 3 <= len(term) <= 50:
            concepts.add(term)
    
    # Find words after colons or dashes (definitions)
    defined_terms = re.findall(r'(?:^|\n)\s*[-–—]\s*([A-Za-z][A-Za-z\s]+?)(?=:|\n)', text)
    for term in defined_terms:
        term = term.strip()
        if 3 <= len(term) <= 50:
            concepts.add(term)
    
    return list(concepts)[:20]  # Return top 20 concepts